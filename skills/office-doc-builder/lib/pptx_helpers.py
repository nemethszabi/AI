"""
Reusable python-pptx helpers - generalized so future slide-deck-generating scripts don't re-derive the
same layout/placeholder boilerplate each time (this framework previously had a one-off script for a
specific slide deck that was deleted once that task finished - these helpers are the reusable capability
that script's own logic should have been split out into, kept separate from the one-time content).

Requires: python-pptx (pip install python-pptx). Note the import name is "pptx".
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN


def new_presentation(widescreen=True):
    """Creates a blank presentation. widescreen=True sets 16:9 (13.33x7.5in); False keeps the default
    4:3 template size."""
    prs = Presentation()
    if widescreen:
        prs.slide_width = Inches(13.333)
        prs.slide_height = Inches(7.5)
    return prs


def add_title_slide(prs, title, subtitle=""):
    """Uses the built-in Title Slide layout (layout index 0 in the default template)."""
    layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    if subtitle and len(slide.placeholders) > 1:
        slide.placeholders[1].text = subtitle
    return slide


def add_content_slide(prs, title, bullets):
    """Uses the built-in Title+Content layout (layout index 1). bullets: list of strings, one per
    top-level bullet point."""
    layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for i, text in enumerate(bullets):
        p = body.paragraphs[0] if i == 0 else body.add_paragraph()
        p.text = text
    return slide


def add_image_slide(prs, image_path, title=None, left_in=1.0, top_in=1.5, width_in=None):
    """Blank layout (index 6 in the default template) with an optional title textbox and one image.
    width_in: if given, scales the image to this width (height auto-scales); if omitted, uses the
    image's native size placed at (left_in, top_in)."""
    layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(layout)
    if title:
        box = slide.shapes.add_textbox(Inches(0.5), Inches(0.3), Inches(12), Inches(1))
        tf = box.text_frame
        tf.text = title
        tf.paragraphs[0].font.size = Pt(28)
        tf.paragraphs[0].font.bold = True
    if width_in:
        slide.shapes.add_picture(image_path, Inches(left_in), Inches(top_in), width=Inches(width_in))
    else:
        slide.shapes.add_picture(image_path, Inches(left_in), Inches(top_in))
    return slide


def set_slide_notes(slide, text):
    slide.notes_slide.notes_text_frame.text = text
